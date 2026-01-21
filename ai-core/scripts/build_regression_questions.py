#!/usr/bin/env python3
"""
Extract test questions from colleague's test documents and build regression CSV.

Inputs:
- Test Questions for Chatbot - Answer Log.docx
- ChatBot Testing Results.pptx

Output:
- test_data/regression_questions.csv
"""

import os
import csv
import re
from pathlib import Path
from typing import List, Set

# Try importing python-docx and python-pptx
try:
    from docx import Document
    DOCX_AVAILABLE = True
except ImportError:
    DOCX_AVAILABLE = False
    print("⚠️ python-docx not available. Install with: pip install python-docx")

try:
    from pptx import Presentation
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False
    print("⚠️ python-pptx not available. Install with: pip install python-pptx")


def is_valid_question(text: str) -> bool:
    """Check if text is a valid user question (not metadata/noise)."""
    # Skip empty or very short
    if not text or len(text) < 5:
        return False
    
    # Skip lines that are clearly metadata
    skip_patterns = [
        r'^\d+[a-z]?\.',  # Numbered items like "1e."
        r'^\-{3,}',  # Separator lines
        r'@miamioh\.edu',  # Email addresses
        r'https?://',  # URLs
        r'^\d+ seconds',  # Timing notes
        r'^[A-Z][A-Z\s]+$',  # ALL CAPS headers
        r'^\d+[a-z]\.',  # "2a." "3b." style
        r'Email:',  # Email labels
        r'Phone:',  # Phone labels
        r'•',  # Bullet point separators
    ]
    
    for pattern in skip_patterns:
        if re.search(pattern, text):
            return False
    
    # Must look like a question
    is_question = any([
        text.endswith("?"),
        text.lower().startswith(("can i", "how do i", "how can", "what", "where", "when", "who", "why", 
                                  "is there", "do you", "does the", "i need", "help me", "find", "tell me",
                                  "show me", "get me", "give me")),
    ])
    
    return is_question


def extract_questions_from_docx(docx_path: str) -> List[str]:
    """Extract questions from Word document."""
    if not DOCX_AVAILABLE:
        print(f"⚠️ Skipping {docx_path} - python-docx not installed")
        return []
    
    questions = []
    
    try:
        doc = Document(docx_path)
        
        for para in doc.paragraphs:
            text = para.text.strip()
            
            if not text:
                continue
            
            # Remove timing notes (e.g., "- 14 seconds")
            text = re.sub(r'\s*\-\s*\d+\s*seconds?\s*$', '', text)
            text = text.strip()
            
            if is_valid_question(text):
                questions.append(text)
        
        print(f"✅ Extracted {len(questions)} questions from {Path(docx_path).name}")
        
    except Exception as e:
        print(f"⚠️ Error reading {docx_path}: {str(e)}")
    
    return questions


def extract_questions_from_pptx(pptx_path: str) -> List[str]:
    """Extract questions from PowerPoint presentation."""
    if not PPTX_AVAILABLE:
        print(f"⚠️ Skipping {pptx_path} - python-pptx not installed")
        return []
    
    questions = []
    
    try:
        prs = Presentation(pptx_path)
        
        for slide in prs.slides:
            for shape in slide.shapes:
                if not hasattr(shape, "text"):
                    continue
                
                text = shape.text.strip()
                
                if not text:
                    continue
                
                # Remove timing notes
                text = re.sub(r'\s*\-\s*\d+\s*seconds?\s*$', '', text)
                text = text.strip()
                
                if is_valid_question(text):
                    questions.append(text)
        
        print(f"✅ Extracted {len(questions)} questions from {Path(pptx_path).name}")
        
    except Exception as e:
        print(f"⚠️ Error reading {pptx_path}: {str(e)}")
    
    return questions


def clean_question(question: str) -> str:
    """Clean and normalize a question string."""
    # Remove bullet points and numbering
    question = re.sub(r'^\s*[\-\*•●○]\s*', '', question)
    question = re.sub(r'^\s*\d+[\.\)]\s*', '', question)
    
    # Remove extra whitespace
    question = ' '.join(question.split())
    
    # Remove trailing punctuation except ? and .
    question = question.rstrip(',:;')
    
    return question


def main():
    """Main extraction logic."""
    print("=" * 80)
    print("Building Regression Test Suite from Colleague Documents")
    print("=" * 80)
    print()
    
    # Check if dependencies are available
    if not DOCX_AVAILABLE and not PPTX_AVAILABLE:
        print("❌ Error: Neither python-docx nor python-pptx is installed")
        print("Install with: pip install python-docx python-pptx")
        return 1
    
    # Define paths
    repo_root = Path(__file__).parent.parent.parent  # chatbot/
    docx_path = repo_root / "Test Questions for Chatbot - Answer Log.docx"
    pptx_path = repo_root / "ChatBot Testing Results.pptx"
    output_dir = Path(__file__).parent.parent / "test_data"
    output_csv = output_dir / "regression_questions.csv"
    
    # Ensure output directory exists
    output_dir.mkdir(exist_ok=True)
    
    # Extract questions from both sources
    all_questions: List[str] = []
    
    if docx_path.exists():
        print(f"📄 Reading: {docx_path}")
        all_questions.extend(extract_questions_from_docx(str(docx_path)))
    else:
        print(f"⚠️ Not found: {docx_path}")
    
    if pptx_path.exists():
        print(f"📊 Reading: {pptx_path}")
        all_questions.extend(extract_questions_from_pptx(str(pptx_path)))
    else:
        print(f"⚠️ Not found: {pptx_path}")
    
    print()
    print(f"📊 Total questions extracted: {len(all_questions)}")
    
    # Clean and deduplicate
    cleaned_questions: Set[str] = set()
    for q in all_questions:
        cleaned = clean_question(q)
        if cleaned and len(cleaned) >= 5:  # Only keep non-empty questions
            cleaned_questions.add(cleaned)
    
    print(f"📊 After cleaning and deduplication: {len(cleaned_questions)}")
    print()
    
    # Sort for consistency
    sorted_questions = sorted(cleaned_questions)
    
    # Write to CSV
    with open(output_csv, 'w', newline='', encoding='utf-8') as f:
        writer = csv.writer(f)
        writer.writerow(['question'])  # Header
        for question in sorted_questions:
            writer.writerow([question])
    
    print(f"✅ Saved to: {output_csv}")
    print(f"✅ Total unique questions: {len(sorted_questions)}")
    print()
    
    # Show sample questions
    print("📝 Sample questions (first 10):")
    for i, q in enumerate(sorted_questions[:10], 1):
        print(f"   {i}. {q}")
    
    print()
    print("=" * 80)
    print("✅ Regression question set ready!")
    print("=" * 80)
    
    return 0


if __name__ == "__main__":
    exit(main())
